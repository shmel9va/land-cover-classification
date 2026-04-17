"""
Generate Land Cover Classification presentation (PPTX) + Speaker Notes (TXT).
Run: py make_presentation.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
import os, copy
from lxml import etree

# ── colours ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x2B, 0x55)   # dark navy – header backgrounds
BLUE   = RGBColor(0x1F, 0x6F, 0xC3)   # medium blue – accents
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # almost-black body text
LGREY  = RGBColor(0xF0, 0xF4, 0xF8)   # light grey slide background
ACCENT = RGBColor(0xE8, 0x8C, 0x00)   # amber – highlight numbers

# ── slide size: 16:9 widescreen ───────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.50)

FIGURES = "results/figures"

def make_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

# ─────────────────────────────────────────────────────────────────────────────
# Helper functions
# ─────────────────────────────────────────────────────────────────────────────

def blank_layout(prs):
    return prs.slide_layouts[6]   # completely blank

def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, italic=False,
                color=DARK, align=PP_ALIGN.LEFT,
                wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_para(tf, text, font_size=18, bold=False, italic=False,
             color=DARK, align=PP_ALIGN.LEFT, bullet=False,
             space_before=Pt(4), font_name="Calibri", level=0):
    """Append a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_before = space_before
    if bullet:
        p._p.get_or_add_pPr().set(qn('a:buChar'), None)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def header_bar(slide, title_text, subtitle_text=""):
    """Dark navy top bar with slide title."""
    add_rect(slide, 0, 0, 13.33, 1.15, NAVY)
    add_textbox(slide, title_text,
                0.25, 0.08, 12.5, 0.75,
                font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle_text:
        add_textbox(slide, subtitle_text,
                    0.25, 0.82, 12.5, 0.30,
                    font_size=13, bold=False, color=RGBColor(0xBB,0xCC,0xDD),
                    align=PP_ALIGN.LEFT)

def accent_line(slide):
    """Thin amber line below header."""
    bar = slide.shapes.add_shape(1, Inches(0), Inches(1.15), Inches(13.33), Pt(3))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

def slide_bg(slide):
    add_rect(slide, 0, 0, 13.33, 7.5, LGREY)

def footer(slide, part_num, total_slides=22):
    """Slide footer with part label."""
    add_rect(slide, 0, 7.1, 13.33, 0.4, NAVY)
    add_textbox(slide,
                "Land Cover Classification · SPbPU · 2024",
                0.2, 7.1, 9, 0.38,
                font_size=10, color=RGBColor(0xAA,0xBB,0xCC))

def img(slide, path, x, y, w=None, h=None):
    if not os.path.exists(path):
        return
    if w and h:
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    elif w:
        slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))
    elif h:
        slide.shapes.add_picture(path, Inches(x), Inches(y), height=Inches(h))
    else:
        slide.shapes.add_picture(path, Inches(x), Inches(y))

def bullet_frame(slide, items, x, y, w, h,
                 font_size=18, color=DARK, title=None, title_color=NAVY):
    """Add a textbox with bullet list."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.name = "Calibri"
        run.font.size = Pt(font_size + 2)
        run.font.bold = True
        run.font.color.rgb = title_color
        first = False
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = ("▸  " if not item.startswith("   ") else "") + item
        run.font.name = "Calibri"
        run.font.size = Pt(font_size)
        run.font.bold = False
        run.font.color.rgb = color
    return txBox

def kpi_box(slide, label, value, x, y, w=2.6, h=1.1,
            label_color=WHITE, value_color=ACCENT, bg=NAVY):
    add_rect(slide, x, y, w, h, bg)
    add_textbox(slide, value, x+0.08, y+0.05, w-0.16, 0.55,
                font_size=26, bold=True, color=value_color, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, x+0.08, y+0.6, w-0.16, 0.45,
                font_size=12, bold=False, color=label_color, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDES
# ─────────────────────────────────────────────────────────────────────────────

def slide_01_title(prs):
    """Part 1 — Title"""
    sl = prs.slides.add_slide(blank_layout(prs))
    # full-bleed navy background
    add_rect(sl, 0, 0, 13.33, 7.5, NAVY)
    # white content area
    add_rect(sl, 0.4, 1.0, 12.53, 5.6, WHITE)
    # amber left stripe
    add_rect(sl, 0.4, 1.0, 0.12, 5.6, ACCENT)

    add_textbox(sl,
        "Land Cover Classification from Satellite Imagery",
        0.65, 1.2, 11.8, 1.2,
        font_size=36, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    add_textbox(sl,
        "Using CNN and U-Net Architectures",
        0.65, 2.35, 11.8, 0.7,
        font_size=30, bold=True, color=BLUE, align=PP_ALIGN.LEFT)

    # divider line
    bar = sl.shapes.add_shape(1, Inches(0.65), Inches(3.1), Inches(11.5), Pt(2))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

    add_textbox(sl,
        "Daniil Khrestyanovskiy   ·   Mariia Shmeleva",
        0.65, 3.25, 11.8, 0.5,
        font_size=20, bold=False, color=DARK, align=PP_ALIGN.LEFT)
    add_textbox(sl,
        "Peter the Great St. Petersburg Polytechnic University\n"
        "Institute of Cybersecurity and Computer Science",
        0.65, 3.85, 11.8, 0.8,
        font_size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT)
    add_textbox(sl,
        "Supervisor: Mohamad Khalil",
        0.65, 4.75, 11.8, 0.5,
        font_size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT)
    add_textbox(sl,
        "2024",
        0.65, 5.4, 11.8, 0.5,
        font_size=14, bold=False, color=RGBColor(0x77,0x88,0x99), align=PP_ALIGN.LEFT)
    return sl

def slide_02_outline(prs):
    """Part 2 — Outline"""
    sl = prs.slides.add_slide(blank_layout(prs))
    slide_bg(sl); header_bar(sl, "Outline"); accent_line(sl); footer(sl, 2)
    items = [
        "1.  Introduction & Problem Statement",
        "2.  State of the Art & Motivation",
        "3.  EuroSAT Dataset",
        "4.  Data Preparation",
        "5.  Reading the Dataset (code)",
        "6.  Model Architectures — SimpleCNN & U-Net Classifier",
        "7.  Model Characteristics",
        "8.  Training Techniques & Hyperparameters",
        "9.  Testing Protocol & Evaluation Metrics",
        "10. Results & Hyperparameter Analysis",
        "11. Conclusion",
        "12. References",
    ]
    txBox = sl.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(11.5), Inches(5.8))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run(); run.text = item
        run.font.name = "Calibri"; run.font.size = Pt(18)
        run.font.color.rgb = DARK if i % 2 == 0 else BLUE
    return sl

def slide_03_aim(prs):
    """Part 3a — Aim of the Work"""
    sl = prs.slides.add_slide(blank_layout(prs))
    slide_bg(sl)
    header_bar(sl, "State of the Art — Aim of the Work",
               "Why do we need automatic land cover classification?")
    accent_line(sl); footer(sl, 3)

    # left column – problem
    add_rect(sl, 0.3, 1.3, 5.9, 5.5, WHITE)
    add_textbox(sl, "The Problem", 0.5, 1.4, 5.5, 0.45,
                font_size=18, bold=True, color=NAVY)
    items_l = [
        "Cities expand rapidly; land use changes constantly",
        "Decision-makers need up-to-date, accurate land maps",
        "Sentinel-2 revisit time: every 5 days over Europe",
        "Volume of imagery: impossible to label manually",
        "Scale demands automation",
    ]
    bullet_frame(sl, items_l, 0.5, 1.9, 5.5, 4.7, font_size=17)

    # right column – objective
    add_rect(sl, 6.5, 1.3, 6.5, 5.5, NAVY)
    add_textbox(sl, "Project Objective", 6.7, 1.4, 6.1, 0.45,
                font_size=18, bold=True, color=ACCENT)
    goals = [
        "Design & train two deep learning architectures",
        "SimpleCNN — lightweight 3-block baseline",
        "U-Net Classifier — encoder-decoder with skip connections",
        "Classify 10 land cover types from satellite RGB images",
        "Compare models on identical data split & metrics",
        "Ensure full reproducibility (fixed seed, Docker)",
    ]
    txBox = sl.shapes.add_textbox(Inches(6.7), Inches(1.95), Inches(6.1), Inches(4.6))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, g in enumerate(goals):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run(); run.text = "▸  " + g
        run.font.name = "Calibri"; run.font.size = Pt(17)
        run.font.color.rgb = WHITE
    return sl

def slide_04_why(prs):
    """Part 3b — Why / Related Work"""
    sl = prs.slides.add_slide(blank_layout(prs))
    slide_bg(sl)
    header_bar(sl, "State of the Art — Related Work & Limitations",
               "What exists and what gap do we fill?")
    accent_line(sl); footer(sl, 4)

    # related work table
    add_textbox(sl, "Key Prior Works", 0.4, 1.3, 6.0, 0.4,
                font_size=18, bold=True, color=NAVY)
    works = [
        ("Helber et al. (2019)", "Introduced EuroSAT; ResNet → 98 % accuracy"),
        ("Ronneberger et al. (2015)", "Original U-Net for biomedical segmentation"),
        ("He et al. (2016)", "ResNet — residual connections for very deep nets"),
        ("Maggiori et al. (2017)", "FCN for pixel-level aerial image labeling"),
        ("Li et al. (2018)", "Attention mechanisms for RS scene classification"),
    ]
    y = 1.75
    for auth, desc in works:
        add_rect(sl, 0.4, y, 2.5, 0.55, NAVY)
        add_textbox(sl, auth, 0.5, y+0.05, 2.3, 0.45, font_size=12, bold=True, color=WHITE)
        add_rect(sl, 2.9, y, 9.8, 0.55, WHITE)
        add_textbox(sl, desc, 3.0, y+0.05, 9.5, 0.45, font_size=14, color=DARK)
        y += 0.65

    # limitations
    add_textbox(sl, "Identified Limitations → Our Contribution", 0.4, 4.7, 12.5, 0.4,
                font_size=18, bold=True, color=NAVY)
    lims = [
        "Deep pre-trained models (ImageNet domain gap) → train from scratch",
        "Multispectral methods need extra sensors → use standard RGB only",
        "Most benchmarks: pixel segmentation → we do scene-level classification",
        "High computational cost → lightweight + efficient U-Net adaptation",
    ]
    bullet_frame(sl, lims, 0.4, 5.15, 12.5, 1.8, font_size=16)
    return sl

def slide_05_dataset_overview(prs):
    """Part 4a — EuroSAT Dataset Overview"""
    sl = prs.slides.add_slide(blank_layout(prs))
    slide_bg(sl)
    header_bar(sl, "EuroSAT Dataset — Overview",
               "eurosat-api.com  ·  Sentinel-2 Satellite (ESA Copernicus)")
    accent_line(sl); footer(sl, 5)

    # KPI boxes
    kpi_box(sl, "Total Images", "27,000", 0.3, 1.35)
    kpi_box(sl, "Image Size", "64 × 64 px", 3.2, 1.35)
    kpi_box(sl, "Land Classes", "10", 6.1, 1.35)
    kpi_box(sl, "Spectral Bands", "13 (RGB used)", 9.0, 1.35)

    # description
    add_rect(sl, 0.3, 2.65, 12.7, 4.1, WHITE)
    txBox = sl.shapes.add_textbox(Inches(0.5), Inches(2.75), Inches(12.3), Inches(3.9))
    tf = txBox.text_frame; tf.word_wrap = True
    lines = [
        ("Dataset source:", "EuroSAT — derived from the Sentinel-2 satellite constellation (ESA Copernicus programme)"),
        ("Spatial resolution:", "10 – 60 m per pixel  |  Revisit time: every 5 days over Europe"),
        ("Spectral bands:", "13 bands available; we use RGB (3 bands) for direct comparability with published benchmarks"),
        ("Normalisation:", "ImageNet statistics: mean [0.485, 0.456, 0.406]  ·  std [0.229, 0.224, 0.225]"),
        ("Class balance:", "AnnualCrop, Forest, Residential: 3,000 samples each  |  Pasture: 2,000  |  others: 2,500"),
        ("Link:", "https://github.com/phelber/EuroSAT"),
    ]
    first = True
    for label, val in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; p.space_before = Pt(6)
        r1 = p.add_run(); r1.text = label + "  "
        r1.font.bold = True; r1.font.color.rgb = NAVY; r1.font.size = Pt(16); r1.font.name = "Calibri"
        r2 = p.add_run(); r2.text = val
        r2.font.color.rgb = DARK; r2.font.size = Pt(16); r2.font.name = "Calibri"
    return sl

def slide_06_classes(prs):
    """Part 4b — 10 Land Cover Classes"""
    sl = prs.slides.add_slide(blank_layout(prs))
    slide_bg(sl)
    header_bar(sl, "EuroSAT Dataset — 10 Land Cover Classes",
               "Each 64×64 patch is labelled with exactly one class")
    accent_line(sl); footer(sl, 6)

    classes = [
        ("AnnualCrop",           "3,000", "Seasonal agricultural fields"),
        ("Forest",               "3,000", "Dense tree cover"),
        ("HerbaceousVegetation", "3,000", "Grassland & meadows"),
        ("Highway",              "2,500", "Major roads & motorways"),
        ("Industrial",           "2,500", "Factories, warehouses"),
        ("Pasture",              "2,000", "Grazing land (smallest class)"),
        ("PermanentCrop",        "2,500", "Orchards & vineyards"),
        ("Residential",          "3,000", "Urban housing areas"),
        ("River",                "2,500", "Waterways & streams"),
        ("SeaLake",              "3,000", "Large open water bodies"),
    ]
    # two columns of 5
    for i, (name, count, desc) in enumerate(classes):
        col = i // 5
        row = i % 5
        x = 0.3 + col * 6.6
        y = 1.35 + row * 1.1
        add_rect(sl, x, y, 6.2, 1.0, WHITE)
        add_rect(sl, x, y, 0.1, 1.0, BLUE)
        add_textbox(sl, name, x+0.2, y+0.05, 3.5, 0.42,
                    font_size=15, bold=True, color=NAVY)
        add_textbox(sl, desc, x+0.2, y+0.48, 3.5, 0.45,
                    font_size=13, color=DARK)
        add_textbox(sl, count, x+3.8, y+0.22, 2.2, 0.45,
                    font_size=20, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
    add_textbox(sl, "images", 6.0, 6.9, 1.5, 0.3,
                font_size=10, color=DARK)
    return sl

def slide_07_samples(prs):
    """Part 4c — Sample Images from Dataset"""
    sl = prs.slides.add_slide(blank_layout(prs))
    slide_bg(sl)
    header_bar(sl, "EuroSAT Dataset — Sample Patches",
               "One representative 64×64 RGB patch per class")
    accent_line(sl); footer(sl, 7)
    fig = os.path.join(FIGURES, "fig2_samples.png")
    img(sl, fig, 0.6, 1.25, w=12.1)
    return sl

def slide_08_prepare_folders(prs):
    """Part 5a — Preparing Dataset: Folder Structure"""
    sl = prs.slides.add_slide(blank_layout(prs))
    slide_bg(sl)
    header_bar(sl, "Preparing the Dataset — Folder Structure",
               "torchvision.datasets.ImageFolder-compatible layout")
    accent_line(sl); footer(sl, 8)

    # folder tree on left
    add_rect(sl, 0.3, 1.3, 6.0, 5.8, RGBColor(0x1E,0x1E,0x2E))
    tree = (
        "EuroSAT/\n"
        "└── 2750/\n"
        "    ├── AnnualCrop/\n"
        "    │   ├── AnnualCrop_00001.jpg\n"
        "    │   └── ...  (3 000 images)\n"
        "    ├── Forest/\n"
        "    │   └── ...  (3 000 images)\n"
        "    ├── HerbaceousVegetation/\n"
        "    ├── Highway/\n"
        "    ├── Industrial/\n"
        "    ├── Pasture/\n"
        "    ├── PermanentCrop/\n"
        "    ├── Residential/\n"
        "    ├── River/\n"
        "    └── SeaLake/\n"
    )
    add_textbox(sl, tree, 0.5, 1.4, 5.6, 5.5,
                font_size=13, color=RGBColor(0xA0,0xFF,0xA0),
                font_name="Courier New")

    # explanation on right
    add_rect(sl, 6.6, 1.3, 6.4, 5.8, WHITE)
    expl = [
        "ImageFolder convention",
        "Each sub-directory name becomes the class label.",
        "The folder index determines the integer class ID (0-9 alphabetically).",
        "",
        "How PyTorch reads it",
        "datasets.ImageFolder(root='EuroSAT/2750')",
        "→ discovers 10 classes automatically",
        "→ returns (image_tensor, class_idx) pairs",
        "",
        "Key benefit",
        "No manual CSV or annotation file needed.",
        "Adding a new class = adding a new sub-folder.",
    ]
    txBox = sl.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.6))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(expl):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        bold = line in ("ImageFolder convention", "How PyTorch reads it", "Key benefit")
        color = NAVY if bold else (BLUE if line.startswith("datasets") or line.startswith("→") else DARK)
        run = p.add_run(); run.text = line
        run.font.name = "Calibri" if not line.startswith("datasets") else "Courier New"
        run.font.size = Pt(14 if not bold else 15)
        run.font.bold = bold
        run.font.color.rgb = color
    return sl

def slide_09_prepare_split(prs):
    """Part 5b — Train / Val / Test Split"""
    sl = prs.slides.add_slide(blank_layout(prs))
    slide_bg(sl)
    header_bar(sl, "Preparing the Dataset — Train / Val / Test Split",
               "Fixed random seed = 42 for full reproducibility")
    accent_line(sl); footer(sl, 9)

    # proportion bar
    bar_x, bar_y, bar_w, bar_h = 0.5, 1.35, 12.33, 0.9
    add_rect(sl, bar_x,               bar_y, bar_w*0.70, bar_h, NAVY)
    add_rect(sl, bar_x+bar_w*0.70,    bar_y, bar_w*0.15, bar_h, BLUE)
    add_rect(sl, bar_x+bar_w*0.85,    bar_y, bar_w*0.15, bar_h, ACCENT)

    add_textbox(sl, "TRAIN  70 %   18,900 images",
                bar_x+0.15, bar_y+0.18, bar_w*0.70-0.3, 0.55,
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, "VAL\n15 %\n4,050",
                bar_x+bar_w*0.70+0.05, bar_y+0.05, bar_w*0.15-0.1, 0.8,
                font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, "TEST\n15 %\n4,050",
                bar_x+bar_w*0.85+0.05, bar_y+0.05, bar_w*0.15-0.1, 0.8,
                font_size=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # details
    details = [
        ("Split method", "torch.utils.data.random_split"),
        ("Random seed", "torch.Generator().manual_seed(42)"),
        ("Training set", "Gradient updates only — never inspected by hand"),
        ("Validation set", "One purpose: select the best epoch (lowest val_loss)"),
        ("Test set", "Touched exactly once, at the end — reports final metrics"),
        ("Leakage prevention", "No hyperparameter tuning uses test set information"),
    ]
    y = 2.6
    for label, val in details:
        add_rect(sl, 0.5, y, 3.5, 0.62, NAVY)
        add_textbox(sl, label, 0.65, y+0.08, 3.2, 0.48,
                    font_size=14, bold=True, color=WHITE)
        add_rect(sl, 4.0, y, 9.0, 0.62, WHITE)
        add_textbox(sl, val, 4.15, y+0.08, 8.7, 0.48,
                    font_size=14, color=DARK)
        y += 0.72
    return sl

def slide_10_reading(prs):
    """Part 6 — Reading the Dataset (code)"""
    sl = prs.slides.add_slide(blank_layout(prs))
    slide_bg(sl)
    header_bar(sl, "Reading the Dataset — data_utils.py",
               "Transforms · DataLoaders · Class discovery")
    accent_line(sl); footer(sl, 10)

    # code block
    add_rect(sl, 0.3, 1.3, 12.7, 4.3, RGBColor(0x1E,0x1E,0x2E))
    code = (
        "from torchvision import datasets, transforms\n"
        "from torch.utils.data import DataLoader, random_split\n\n"
        "transform = transforms.Compose([\n"
        "    transforms.ToTensor(),\n"
        "    transforms.Normalize(mean=[0.485, 0.456, 0.406],\n"
        "                         std= [0.229, 0.224, 0.225])\n"
        "])\n\n"
        "full_dataset = datasets.ImageFolder(root='EuroSAT/2750', transform=transform)\n"
        "# classes: ['AnnualCrop', 'Forest', ..., 'SeaLake']  (10 total)\n\n"
        "train_dataset, val_dataset, test_dataset = random_split(\n"
        "    full_dataset, [18900, 4050, 4050],\n"
        "    generator=torch.Generator().manual_seed(42)\n"
        ")\n"
        "train_loader = DataLoader(train_dataset, batch_size=64, shuffle=True)"
    )
    add_textbox(sl, code, 0.5, 1.4, 12.3, 4.1,
                font_size=13, color=RGBColor(0xA0,0xD0,0xFF),
                font_name="Courier New")

    # annotations below code
    notes = [
        "ToTensor()   — converts PIL image [0,255] → float tensor [0.0, 1.0]",
        "Normalize()  — subtracts ImageNet channel mean, divides by channel std",
        "ImageFolder  — auto-discovers classes from sub-directory names",
        "batch_size=64 — good balance between gradient noise and memory",
    ]
    bullet_frame(sl, notes, 0.3, 5.7, 12.7, 1.6, font_size=14)
    return sl

def slide_11_simplecnn(prs):
    """Part 7a — SimpleCNN Architecture"""
    sl = prs.slides.add_slide(blank_layout(prs))
    slide_bg(sl)
    header_bar(sl, "Creating the CNN — SimpleCNN Architecture",
               "Lightweight 3-block convolutional baseline  ·  ~4.5 M parameters")
    accent_line(sl); footer(sl, 11)

    # architecture diagram using boxes
    blocks = [
        ("Input\n3×64×64",     0.3,  2.0, 1.6, 2.8, LGREY, DARK),
        ("Conv Block 1\n3→32 ch\n64→32 px\nBN+ReLU+Pool", 2.2, 1.6, 2.0, 3.5, NAVY, WHITE),
        ("Conv Block 2\n32→64 ch\n32→16 px\nBN+ReLU+Pool", 4.5, 1.6, 2.0, 3.5, NAVY, WHITE),
        ("Conv Block 3\n64→128 ch\n16→8 px\nBN+ReLU+Pool", 6.8, 1.6, 2.0, 3.5, NAVY, WHITE),
        ("Flatten\n8192",      9.1,  2.0, 1.4, 2.8, BLUE, WHITE),
        ("FC 8192→512\nDropout 0.5\nReLU", 10.8, 1.8, 1.5, 3.1, BLUE, WHITE),
        ("FC 512→10\nSoftmax",  12.5, 2.2, 0.7, 2.4, ACCENT, DARK),
    ]
    for label, x, y, w, h, bg, fg in blocks:
        add_rect(sl, x, y, w, h, bg)
        add_textbox(sl, label, x+0.05, y+0.1, w-0.1, h-0.2,
                    font_size=13, bold=False, color=fg, align=PP_ALIGN.CENTER)
    # arrows
    arrs = [(1.9, 3.35), (4.2, 3.35), (6.5, 3.35), (9.05, 3.35), (10.5, 3.35), (12.3, 3.35)]
    for ax, ay in arrs:
        add_textbox(sl, "→", ax, ay, 0.3, 0.4, font_size=18, bold=True, color=NAVY)

    # spatial flow label
    add_textbox(sl, "Spatial: 64 → 32 → 16 → 8     Channels: 3 → 32 → 64 → 128",
                1.5, 5.5, 11.0, 0.5, font_size=14, color=DARK, align=PP_ALIGN.CENTER)

    # kernel / pool info
    info = [
        "3×3 conv, padding=1 (preserves spatial size before pooling)",
        "2×2 MaxPool halves spatial dimensions at each stage",
        "Dropout p=0.5 in FC head prevents overfitting",
    ]
    bullet_frame(sl, info, 0.4, 6.05, 12.5, 1.2, font_size=14)
    return sl

def slide_12_unet(prs):
    """Part 7b — U-Net Classifier Architecture"""
    sl = prs.slides.add_slide(blank_layout(prs))
    slide_bg(sl)
    header_bar(sl, "Creating the CNN — U-Net Classifier Architecture",
               "Encoder-decoder with skip connections  ·  ~8.1 M parameters")
    accent_line(sl); footer(sl, 12)

    # encoder column
    add_rect(sl, 0.3, 1.35, 3.8, 5.8, WHITE)
    add_textbox(sl, "ENCODER", 0.5, 1.38, 3.4, 0.4,
                font_size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    enc = [
        ("Input  3ch  64×64",  1.4),
        ("DoubleConv  64ch",   1.9),
        ("Pool + DC  128ch",   2.55),
        ("Pool + DC  256ch",   3.2),
        ("Pool + DC  512ch  8×8", 3.85),
    ]
    for txt, y in enc:
        add_rect(sl, 0.5, y, 3.3, 0.55, NAVY)
        add_textbox(sl, txt, 0.55, y+0.07, 3.2, 0.42,
                    font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

    # decoder column
    add_rect(sl, 4.5, 1.35, 3.8, 5.8, WHITE)
    add_textbox(sl, "DECODER", 4.7, 1.38, 3.4, 0.4,
                font_size=15, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    dec = [
        ("TransConv 256ch",    1.9),
        ("cat+DC  256ch",      2.55),
        ("TransConv 128ch",    3.2),
        ("cat+DC  128ch",      3.85),
        ("TransConv 64ch",     4.5),
        ("cat+DC  64ch",       5.15),
    ]
    for txt, y in dec:
        add_rect(sl, 4.7, y, 3.3, 0.55, BLUE)
        add_textbox(sl, txt, 4.75, y+0.07, 3.2, 0.42,
                    font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

    # classification head
    add_rect(sl, 8.7, 1.35, 4.3, 5.8, WHITE)
    add_textbox(sl, "CLASSIFICATION HEAD", 8.9, 1.38, 3.9, 0.4,
                font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    head = [
        ("AdaptiveAvgPool2d\n(1×1)",        2.3),
        ("Flatten → 64-dim",                3.2),
        ("Linear 64 → 10",                  4.0),
        ("Class Logits",                    4.8),
    ]
    for txt, y in head:
        add_rect(sl, 8.9, y, 3.9, 0.7, ACCENT)
        add_textbox(sl, txt, 8.95, y+0.05, 3.8, 0.6,
                    font_size=13, color=DARK, align=PP_ALIGN.CENTER)

    # skip connection labels
    for y in [2.1, 2.75, 3.4]:
        add_textbox(sl, "skip ─────▶", 3.85, y+0.05, 0.85, 0.4,
                    font_size=10, color=RGBColor(0xCC,0x44,0x00), bold=True)

    add_textbox(sl,
        "Key adaptation: segmentation head replaced by GlobalAvgPool + Linear → one label per image",
        0.3, 7.05, 13.0, 0.35, font_size=12, color=DARK, align=PP_ALIGN.CENTER)
    return sl

def slide_13_arch_compare(prs):
    """Part 7c — Architecture Comparison (workflow figure)"""
    sl = prs.slides.add_slide(blank_layout(prs))
    slide_bg(sl)
    header_bar(sl, "Architecture Overview — Pipeline",
               "EuroSAT  →  Preprocessing  →  Model  →  Training  →  Evaluation")
    accent_line(sl); footer(sl, 13)
    fig = os.path.join(FIGURES, "fig1_workflow.png")
    img(sl, fig, 0.4, 1.25, w=12.5)
    return sl

def slide_14_cnn_chars(prs):
    """Part 8a — SimpleCNN Characteristics"""
    sl = prs.slides.add_slide(blank_layout(prs))
    slide_bg(sl)
    header_bar(sl, "Model Characteristics — SimpleCNN",
               "Layer-by-layer parameter breakdown")
    accent_line(sl); footer(sl, 14)

    rows = [
        ("Layer",             "Output shape",   "Parameters"),
        ("Conv2d  3→32",      "32×64×64",       "896"),
        ("BatchNorm2d 32",    "32×64×64",        "64"),
        ("MaxPool2d  2×2",    "32×32×32",       "0"),
        ("Conv2d  32→64",     "64×32×32",       "18,496"),
        ("BatchNorm2d 64",    "64×32×32",        "128"),
        ("MaxPool2d  2×2",    "64×16×16",       "0"),
        ("Conv2d  64→128",    "128×16×16",      "73,856"),
        ("BatchNorm2d 128",   "128×16×16",       "256"),
        ("MaxPool2d  2×2",    "128×8×8",        "0"),
        ("Flatten",           "8,192",          "0"),
        ("Linear 8192→512",   "512",            "4,194,816"),
        ("Dropout 0.5",       "512",            "0"),
        ("Linear 512→10",     "10",             "5,130"),
        ("TOTAL",             "",               "~4.5 M"),
    ]
    col_w = [4.2, 3.6, 3.6]
    col_x = [0.3, 4.7, 8.5]
    y = 1.35
    for i, row in enumerate(rows):
        bg = NAVY if i == 0 else (LGREY if i % 2 == 0 else WHITE)
        fg = WHITE if i == 0 else DARK
        is_total = (i == len(rows)-1)
        if is_total: bg = ACCENT
        h = 0.38
        for j, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
            add_rect(sl, cx, y, cw-0.05, h, bg)
            add_textbox(sl, cell, cx+0.1, y+0.04, cw-0.2, h-0.08,
                        font_size=12, bold=(i==0 or is_total),
                        color=(WHITE if i==0 or is_total else DARK))
        y += h

    kpi_box(sl, "Total Parameters", "4.5 M",  0.3,  6.1, w=3.0)
    kpi_box(sl, "Model Size",       "17.2 MB", 3.5,  6.1, w=3.0)
    kpi_box(sl, "Trained (best epoch)", "11", 6.7,  6.1, w=3.0)
    kpi_box(sl, "Test Accuracy",   "89.23 %",  9.9,  6.1, w=3.0)
    return sl

def slide_15_unet_chars(prs):
    """Part 8b — U-Net Classifier Characteristics"""
    sl = prs.slides.add_slide(blank_layout(prs))
    slide_bg(sl)
    header_bar(sl, "Model Characteristics — U-Net Classifier",
               "Encoder · Decoder · Classification Head breakdown")
    accent_line(sl); footer(sl, 15)

    rows = [
        ("Block",                     "Output shape",  "Approx. params"),
        ("DoubleConv  3→64   (inc)",   "64×64×64",     "~38 K"),
        ("MaxPool + DoubleConv 64→128","128×32×32",    "~221 K"),
        ("MaxPool + DoubleConv 128→256","256×16×16",   "~885 K"),
        ("MaxPool + DoubleConv 256→512","512×8×8",     "~3.5 M"),
        ("ConvTranspose  512→256",     "256×16×16",    "~524 K"),
        ("DoubleConv 512→256  (skip)", "256×16×16",    "~1.2 M"),
        ("ConvTranspose  256→128",     "128×32×32",    "~131 K"),
        ("DoubleConv 256→128  (skip)", "128×32×32",    "~590 K"),
        ("ConvTranspose  128→64",      "64×64×64",     "~32 K"),
        ("DoubleConv 128→64   (skip)", "64×64×64",     "~148 K"),
        ("AdaptiveAvgPool + Linear",   "10",           "650"),
        ("TOTAL",                      "",             "~8.1 M"),
    ]
    col_w = [4.6, 3.4, 3.4]
    col_x = [0.3, 5.1, 8.7]
    y = 1.35
    for i, row in enumerate(rows):
        bg = NAVY if i == 0 else (LGREY if i % 2 == 0 else WHITE)
        is_total = (i == len(rows)-1)
        if is_total: bg = ACCENT
        h = 0.37
        for cell, cx, cw in zip(row, col_x, col_w):
            add_rect(sl, cx, y, cw-0.05, h, bg)
            add_textbox(sl, cell, cx+0.1, y+0.04, cw-0.2, h-0.08,
                        font_size=11, bold=(i==0 or is_total),
                        color=(WHITE if i==0 or is_total else DARK))
        y += h

    kpi_box(sl, "Total Parameters", "8.1 M",  0.3,  6.15, w=3.0)
    kpi_box(sl, "Model Size",       "30.9 MB", 3.5,  6.15, w=3.0)
    kpi_box(sl, "Trained (best epoch)", "24", 6.7,  6.15, w=3.0)
    kpi_box(sl, "Test Accuracy",   "96.10 %",  9.9,  6.15, w=3.0)
    return sl

def slide_16_training(prs):
    """Part 9 — Training Techniques & Hyperparameters"""
    sl = prs.slides.add_slide(blank_layout(prs))
    slide_bg(sl)
    header_bar(sl, "Training Techniques & Hyperparameters",
               "Identical conditions for both models — fair comparison")
    accent_line(sl); footer(sl, 16)

    params = [
        ("Optimizer",    "Adam (lr = 0.001)  —  stable convergence for both models"),
        ("Loss function","CrossEntropyLoss  —  standard multi-class classification"),
        ("Batch size",   "64  —  balances gradient noise vs. GPU memory"),
        ("Max epochs",   "50  —  upper bound; early stopping usually triggers earlier"),
        ("Early stopping","Patience = 5 epochs; min_delta = 0.001 (CNN) / 0.0 (U-Net)"),
        ("Checkpoint",   "Best model saved when val_loss improves; used for final eval"),
        ("Random seed",  "42 for both data split and weight init — full reproducibility"),
        ("Device",       "CUDA GPU if available, else CPU"),
    ]
    y = 1.35
    for label, val in params:
        add_rect(sl, 0.3, y, 3.5, 0.62, NAVY)
        add_textbox(sl, label, 0.45, y+0.09, 3.2, 0.46,
                    font_size=14, bold=True, color=WHITE)
        add_rect(sl, 3.85, y, 9.1, 0.62, WHITE)
        add_textbox(sl, val, 4.0, y+0.09, 8.8, 0.46,
                    font_size=14, color=DARK)
        y += 0.72

    # early stopping benefit note
    add_rect(sl, 0.3, 7.05, 12.7, 0.38, ACCENT)
    add_textbox(sl,
        "Early stopping saved 34 epochs for SimpleCNN · 21 epochs for U-Net  —  "
        "protects against overfitting by selecting best validation checkpoint",
        0.45, 7.06, 12.4, 0.36, font_size=12, bold=True, color=DARK)
    return sl

def slide_17_testing(prs):
    """Part 10 — Testing"""
    sl = prs.slides.add_slide(blank_layout(prs))
    slide_bg(sl)
    header_bar(sl, "Testing Protocol & Evaluation Metrics",
               "4,050 held-out images  ·  never seen during training or tuning")
    accent_line(sl); footer(sl, 17)

    # metrics boxes
    metrics_info = [
        ("Accuracy", "Fraction of correctly classified images\n(all 4,050 test samples)"),
        ("Weighted F1", "Harmonic mean of precision & recall\nweighted by class frequency"),
        ("Weighted\nPrecision", "TP / (TP + FP)  per class\nweighted by class frequency"),
        ("Weighted\nRecall", "TP / (TP + FN)  per class\nweighted by class frequency"),
    ]
    for i, (name, desc) in enumerate(metrics_info):
        x = 0.3 + i * 3.25
        add_rect(sl, x, 1.35, 3.0, 1.8, NAVY)
        add_textbox(sl, name, x+0.1, 1.4, 2.8, 0.55,
                    font_size=17, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(sl, desc, x+0.1, 2.0, 2.8, 1.1,
                    font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

    # why weighted
    add_rect(sl, 0.3, 3.35, 12.7, 1.05, WHITE)
    add_textbox(sl, "Why weighted averaging?", 0.5, 3.38, 4.0, 0.4,
                font_size=16, bold=True, color=NAVY)
    add_textbox(sl,
        "Class sizes differ (Pasture: 2,000 · Forest: 3,000). Weighted metrics reflect the real "
        "class distribution, giving a more realistic picture of deployment performance than macro averaging.",
        0.5, 3.8, 12.3, 0.55, font_size=14, color=DARK)

    # confusion matrix note
    add_rect(sl, 0.3, 4.55, 12.7, 1.0, LGREY)
    add_textbox(sl, "Confusion matrix (10 × 10)", 0.5, 4.58, 5.0, 0.4,
                font_size=16, bold=True, color=NAVY)
    add_textbox(sl,
        "Each cell [i, j] shows how many test images of true class i were predicted as class j.\n"
        "Diagonal = correct predictions. Off-diagonal = errors. Used to identify the hardest class pairs.",
        0.5, 5.0, 12.3, 0.5, font_size=14, color=DARK)

    # best-epoch rule
    add_rect(sl, 0.3, 5.65, 12.7, 0.8, NAVY)
    add_textbox(sl,
        "Checkpoint selection rule: evaluate.py loads the checkpoint with the lowest validation loss "
        "(not the last epoch). Validation loss is a more sensitive signal than accuracy: it degrades "
        "before accuracy does, giving earlier warning of overfitting.",
        0.5, 5.68, 12.3, 0.75, font_size=13, color=WHITE)

    add_rect(sl, 0.3, 6.52, 12.7, 0.45, ACCENT)
    add_textbox(sl,
        "Test set is used EXACTLY ONCE  —  after all hyperparameter choices are final",
        0.5, 6.54, 12.3, 0.4, font_size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    return sl

def slide_18_results_cnn(prs):
    """Part 11a — Results: SimpleCNN"""
    sl = prs.slides.add_slide(blank_layout(prs))
    slide_bg(sl)
    header_bar(sl, "Results — SimpleCNN",
               "Test set · 4,050 images · best checkpoint at epoch 11")
    accent_line(sl); footer(sl, 18)

    kpi_box(sl, "Test Accuracy", "89.23 %", 0.3,  1.35, w=2.8)
    kpi_box(sl, "Weighted F1",   "0.8930",  3.3,  1.35, w=2.8)
    kpi_box(sl, "Precision",     "0.8963",  6.3,  1.35, w=2.8)
    kpi_box(sl, "Recall",        "0.8923",  9.3,  1.35, w=2.8)

    # training curves
    fig = os.path.join(FIGURES, "fig3_cnn_curves.png")
    img(sl, fig, 0.3, 2.6, w=7.8)

    # per-class highlights
    add_rect(sl, 8.3, 2.6, 4.7, 4.25, WHITE)
    add_textbox(sl, "Per-Class F1 Highlights", 8.45, 2.65, 4.4, 0.4,
                font_size=15, bold=True, color=NAVY)
    highlights = [
        ("Best classes",        ""),
        ("Forest",              "0.954"),
        ("SeaLake",             "0.948"),
        ("Residential",         "0.939"),
        ("Hardest classes",     ""),
        ("HerbaceousVegetation","0.821  ← visually similar to Pasture"),
        ("PermanentCrop",       "0.833  ← texture overlap with AnnualCrop"),
        ("River",               "0.836  ← thin strip, confused with Highway"),
    ]
    y = 3.1
    for label, val in highlights:
        if val == "":
            add_textbox(sl, label, 8.45, y, 4.4, 0.35,
                        font_size=13, bold=True, color=BLUE)
        else:
            add_textbox(sl, f"  {label}: {val}", 8.45, y, 4.4, 0.38,
                        font_size=12, color=DARK)
        y += 0.42
    return sl

def slide_19_results_unet(prs):
    """Part 11b — Results: U-Net Classifier"""
    sl = prs.slides.add_slide(blank_layout(prs))
    slide_bg(sl)
    header_bar(sl, "Results — U-Net Classifier",
               "Test set · 4,050 images · best checkpoint at epoch 24")
    accent_line(sl); footer(sl, 19)

    kpi_box(sl, "Test Accuracy", "96.10 %", 0.3,  1.35, w=2.8)
    kpi_box(sl, "Weighted F1",   "0.9610",  3.3,  1.35, w=2.8)
    kpi_box(sl, "Precision",     "0.9613",  6.3,  1.35, w=2.8)
    kpi_box(sl, "Recall",        "0.9610",  9.3,  1.35, w=2.8)

    fig = os.path.join(FIGURES, "fig4_unet_curves.png")
    img(sl, fig, 0.3, 2.6, w=7.8)

    add_rect(sl, 8.3, 2.6, 4.7, 4.25, WHITE)
    add_textbox(sl, "Gains over SimpleCNN", 8.45, 2.65, 4.4, 0.4,
                font_size=15, bold=True, color=NAVY)
    gains = [
        ("River",                "0.836 → 0.951  (+0.115)"),
        ("PermanentCrop",        "0.833 → 0.940  (+0.107)"),
        ("HerbaceousVeg.",       "0.821 → 0.939  (+0.118)"),
        ("Highway",              "0.848 → 0.978  (+0.130)"),
        ("Residential",          "0.939 → 0.986  (+0.047)"),
        ("SeaLake",              "0.948 → 0.979  (+0.031)"),
    ]
    y = 3.1
    for cls, gain in gains:
        add_rect(sl, 8.45, y, 4.4, 0.38, LGREY)
        add_textbox(sl, f"{cls}: {gain}", 8.55, y+0.04, 4.2, 0.32,
                    font_size=12, color=DARK)
        y += 0.45
    add_textbox(sl,
        "Improvement is consistent across ALL classes,\n"
        "not limited to easy ones → systematic architectural gain",
        8.45, 5.75, 4.4, 0.9, font_size=13, color=NAVY, bold=False)
    return sl

def slide_20_comparison(prs):
    """Part 11c — Results Comparison"""
    sl = prs.slides.add_slide(blank_layout(prs))
    slide_bg(sl)
    header_bar(sl, "Results — Model Comparison",
               "SimpleCNN  vs.  U-Net Classifier  ·  same test set  ·  same metrics")
    accent_line(sl); footer(sl, 20)

    # comparison table
    rows = [
        ("Metric",            "SimpleCNN",  "U-Net Classifier", "Δ"),
        ("Test Accuracy",     "89.23 %",    "96.10 %",          "+6.87 pp"),
        ("Weighted F1",       "0.8930",     "0.9610",           "+0.068"),
        ("Weighted Precision","0.8963",     "0.9613",           "+0.065"),
        ("Weighted Recall",   "0.8923",     "0.9610",           "+0.069"),
        ("Best Val Loss",     "0.3105",     "0.1116",           "−64 %"),
        ("Best Epoch",        "11",         "24",               "×2.2"),
        ("Parameters",        "~4.5 M",     "~8.1 M",           "×1.8"),
        ("Model size",        "17.2 MB",    "30.9 MB",          "×1.8"),
        ("Worst-class F1",    "0.821",      "0.939",            "+0.118"),
    ]
    col_x = [0.3, 4.5, 7.8, 11.2]
    col_w = [4.0, 3.1, 3.2, 1.8]
    y = 1.35
    for i, row in enumerate(rows):
        bg = NAVY if i == 0 else (LGREY if i % 2 == 0 else WHITE)
        fg = WHITE if i == 0 else DARK
        h = 0.48
        for cell, cx, cw in zip(row, col_x, col_w):
            add_rect(sl, cx, y, cw-0.05, h, bg)
            delta_color = RGBColor(0x00, 0x88, 0x00) if (i > 0 and "+" in cell) else (
                          RGBColor(0xCC, 0x00, 0x00) if (i > 0 and "−" in cell) else fg)
            add_textbox(sl, cell, cx+0.1, y+0.06, cw-0.2, h-0.12,
                        font_size=12, bold=(i==0),
                        color=delta_color if i > 0 and cx == col_x[3] else fg)
        y += h

    # bar chart image
    fig = os.path.join(FIGURES, "fig6_metrics_comparison.png")
    img(sl, fig, 0.3, 6.1, h=1.25)

    fig2 = os.path.join(FIGURES, "fig5_f1_comparison.png")
    img(sl, fig2, 6.8, 6.1, h=1.25)
    return sl

def slide_21_conclusion(prs):
    """Part 12 — Conclusion"""
    sl = prs.slides.add_slide(blank_layout(prs))
    slide_bg(sl)
    header_bar(sl, "Conclusion", "")
    accent_line(sl); footer(sl, 21)

    left = [
        "SimpleCNN",
        "89.23 % accuracy  ·  F1 = 0.893",
        "Converges in 11 epochs  ·  4.5 M params",
        "Fast, lightweight, practical baseline",
        "",
        "U-Net Classifier",
        "96.10 % accuracy  ·  F1 = 0.961",
        "Converges in 24 epochs  ·  8.1 M params",
        "+6.87 pp over SimpleCNN across all classes",
    ]
    txBox = sl.shapes.add_textbox(Inches(0.4), Inches(1.35), Inches(6.0), Inches(5.5))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(left):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        bold = line in ("SimpleCNN", "U-Net Classifier")
        color = NAVY if bold else DARK
        fs = 17 if bold else 15
        run = p.add_run(); run.text = ("" if bold else "▸  ") + line
        run.font.name = "Calibri"; run.font.size = Pt(fs)
        run.font.bold = bold; run.font.color.rgb = color

    add_rect(sl, 6.7, 1.35, 6.3, 5.5, NAVY)
    right = [
        "Key Lessons",
        "Skip connections + hierarchical reuse = best architectural choice for satellite imagery",
        "U-Net naturally combines local texture and global context",
        "Batch Normalization critical: +3 pp accuracy, faster convergence",
        "Early stopping + fixed seed → reproducible, robust evaluation",
        "",
        "Future Work",
        "Use all 13 Sentinel-2 bands (NIR/SWIR boost vegetation classes)",
        "Data augmentation (flips, rotation) especially for Pasture",
        "Transfer learning: ResNet/EfficientNet backbone fine-tuning",
        "Temporal fusion: multi-date stack for seasonal dynamics",
    ]
    txBox2 = sl.shapes.add_textbox(Inches(6.9), Inches(1.4), Inches(5.9), Inches(5.3))
    tf2 = txBox2.text_frame; tf2.word_wrap = True
    for i, line in enumerate(right):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_before = Pt(4)
        bold = line in ("Key Lessons", "Future Work")
        color = ACCENT if bold else WHITE
        fs = 16 if bold else 13
        run = p.add_run(); run.text = ("" if (bold or line == "") else "▸  ") + line
        run.font.name = "Calibri"; run.font.size = Pt(fs)
        run.font.bold = bold; run.font.color.rgb = color
    return sl

def slide_22_references(prs):
    """Part 13 — References"""
    sl = prs.slides.add_slide(blank_layout(prs))
    slide_bg(sl)
    header_bar(sl, "References", "")
    accent_line(sl); footer(sl, 22)

    refs = [
        '[1]  P. Helber, B. Bischke, A. Dengel, D. Borth. "EuroSAT: A Novel Dataset and Deep Learning'
        ' Benchmark for Land Use and Land Cover Classification." IEEE JSTARS, 2019.',
        '[2]  O. Ronneberger, P. Fischer, T. Brox. "U-Net: Convolutional Networks for Biomedical Image'
        ' Segmentation." MICCAI, 2015.',
        '[3]  K. He, X. Zhang, S. Ren, J. Sun. "Deep Residual Learning for Image Recognition."'
        ' CVPR, 2016.',
        '[4]  E. Maggiori, Y. Tarabalka, G. Charpiat, P. Alliez. "Convolutional Neural Networks for'
        ' Large-Scale Remote-Sensing Image Classification." IEEE TGRS, 2017.',
        '[5]  K. S. Nogueira et al. "Exploiting EfficientNet and Transfer Learning for Land Use'
        ' Classification." Remote Sensing, 2019.',
        '[6]  Y. Li et al. "Deep Learning for Remote Sensing Image Classification: A Survey."'
        ' Wiley WIREs Data Mining, 2018.',
        '[7]  S. Ioffe, C. Szegedy. "Batch Normalization: Accelerating Deep Network Training."'
        ' ICML, 2015.',
        '[8]  D. P. Kingma, J. Ba. "Adam: A Method for Stochastic Optimization." ICLR, 2015.',
    ]
    txBox = sl.shapes.add_textbox(Inches(0.4), Inches(1.35), Inches(12.5), Inches(5.8))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, ref in enumerate(refs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(7)
        run = p.add_run(); run.text = ref
        run.font.name = "Calibri"; run.font.size = Pt(14); run.font.color.rgb = DARK
    return sl

# ─────────────────────────────────────────────────────────────────────────────
# SPEAKER NOTES
# ─────────────────────────────────────────────────────────────────────────────

NOTES = {
    1: """[~1:00]
Good morning, everyone. My name is [your name], and I am presenting our project "Land Cover Classification from Satellite Imagery Using CNN and U-Net Architectures." This work was carried out at Peter the Great St. Petersburg Polytechnic University, Institute of Cybersecurity and Computer Science, under the supervision of Mohamad Khalil. Our goal was to train two deep learning models — a simple baseline CNN and a U-Net-based classifier — that automatically determine what type of land cover appears in a 64×64 satellite image patch from the EuroSAT dataset. The talk will take 20 to 25 minutes, after which we are happy to take questions.""",

    2: """[~0:30]
Let me walk you through what we will cover. We start with the problem statement and motivation, move to related work and our contribution, then describe the EuroSAT dataset in detail. After that we look at data preparation and how it is loaded in code. The bulk of the talk covers the two model architectures, how they were trained, tested, and how the results compare. We close with conclusions and references.""",

    3: """[~1:30]
Why do we need automatic land cover classification? The scale of modern Earth observation makes manual labeling impractical. Sentinel-2 covers the entire European continent every five days — one week of imagery would take a team of analysts months to review by hand. Decision-makers in urban planning, agriculture, and environmental monitoring need frequent, accurate land use maps. Our objective is clear: design and fairly compare two architectures — a lightweight SimpleCNN baseline and a U-Net repurposed for scene-level classification — trained entirely from scratch on EuroSAT RGB imagery, with no ImageNet pre-training.""",

    4: """[~1:30]
Before our work, several key papers defined this landscape. Helber et al. introduced EuroSAT in 2019 and showed fine-tuned ResNet achieves over 98% — our upper-bound reference. U-Net, introduced by Ronneberger et al. for biomedical segmentation, inspired our encoder-decoder architecture. He et al.'s ResNet showed residual connections enable very deep networks, and Li et al. later added attention mechanisms for remote sensing. These works set the benchmark. Our contribution: we address the domain gap of ImageNet pre-training by training from scratch, and we adapt U-Net for scene-level classification instead of pixel-wise segmentation, making it suitable for EuroSAT's patch-labeling task.""",

    5: """[~1:30]
The EuroSAT dataset comes from the Sentinel-2 satellite constellation, part of ESA's Copernicus programme. It contains 27,000 image patches, each 64×64 pixels, across 10 land cover classes. Although Sentinel-2 provides 13 spectral bands at resolutions from 10 to 60 meters per pixel, we use only the RGB version — the three visible bands. This makes our results directly comparable with published benchmarks. Before feeding images into the network, we normalize each channel using ImageNet statistics: mean [0.485, 0.456, 0.406] and standard deviation [0.229, 0.224, 0.225]. This is standard practice for RGB CNNs. The class distribution is slightly imbalanced: Forest, AnnualCrop, and Residential each have 3,000 samples, while Pasture has only 2,000. This mild imbalance is worth keeping in mind when looking at per-class results.""",

    6: """[~1:00]
Here are the 10 land cover classes. Notice the variation in sample counts — Pasture has the fewest images at 2,000, which makes it potentially the hardest class to learn. The classes span a wide range of visual appearances: from dense Forest and open SeaLake — which are visually distinctive — to HerbaceousVegetation and PermanentCrop, which are difficult to separate even for a human expert at this resolution. This diversity is what makes EuroSAT a meaningful benchmark.""",

    7: """[~0:45]
Here you can see one sample patch for each of the 10 classes. Notice how Forest looks dark and textured, SeaLake is uniformly dark blue, and Residential shows the regular grid of streets and rooftops. The hardest pairs are immediately apparent: HerbaceousVegetation and Pasture look almost identical; PermanentCrop and AnnualCrop share similar textures. At 64×64 pixels and 10 meters per pixel, these distinctions are genuinely challenging.""",

    8: """[~1:00]
The dataset lives on disk in an ImageFolder-compatible layout: one sub-directory per class, images inside. This is the format expected by torchvision's ImageFolder class. The class integer ID is assigned alphabetically — AnnualCrop gets index 0, Forest gets 1, and so on up to SeaLake at index 9. No annotation files, no CSV — the directory structure is the label. The key benefit is simplicity: adding a new class means adding a new folder.""",

    9: """[~1:00]
We split the 27,000 images 70 / 15 / 15 using torch.utils.data.random_split with a fixed generator seed of 42. That gives 18,900 training images, 4,050 for validation, and 4,050 for the test set. The roles are strictly separated. Training set: gradient updates only. Validation set: one purpose — choose which epoch's checkpoint to save. Test set: touched exactly once at the end to report final metrics. This strict separation is the foundation of our unbiased evaluation.""",

    10: """[~1:00]
Here is the actual code from data_utils.py. Notice the transform pipeline: ToTensor converts a PIL image from uint8 [0,255] to float [0.0,1.0], then Normalize subtracts the ImageNet per-channel mean and divides by the standard deviation. ImageFolder discovers all 10 classes automatically from the directory names. random_split with the fixed seed produces the same three disjoint subsets on every run — this is what makes our experiments reproducible. Batch size 64 feeds the DataLoader: large enough for stable gradient estimates, small enough to fit comfortably in GPU memory.""",

    11: """[~1:45]
SimpleCNN is our baseline. It has three convolutional blocks, each with a 3×3 convolution with padding=1 — which preserves spatial size before pooling — followed by Batch Normalization, ReLU activation, and 2×2 MaxPool that halves spatial dimensions. The channel depth grows: 3 → 32 → 64 → 128. After three pools, the 8×8 feature map is flattened to a vector of 8,192 elements. Two fully connected layers follow: 8192→512 with Dropout 0.5 to prevent overfitting, then 512→10 for the class logits. Total: ~4.5 million parameters. Batch Normalization proved critical — removing it reduced accuracy by ~3 percentage points in our preliminary experiments.""",

    12: """[~1:45]
The U-Net Classifier has an encoder path of four DoubleConv blocks — two 3×3 convolutions per stage, each with BatchNorm and ReLU — interspersed with MaxPool operations. Channels grow: 3→64→128→256→512, while spatial size shrinks from 64×64 down to 8×8. The decoder mirrors this with transposed convolutions that upsample back to full resolution. The key innovation is skip connections: before each DoubleConv in the decoder, we concatenate the upsampled feature map with the matching encoder feature map, preserving fine spatial detail that would otherwise be lost during downsampling. Our adaptation for classification: instead of the original 1×1 segmentation head, we attach an AdaptiveAvgPool that collapses the spatial map to a single 64-dimensional vector, followed by a Linear layer to 10 classes. Total: ~8.1 million parameters.""",

    13: """[~1:00]
This workflow diagram summarizes the full pipeline. On the left: the EuroSAT dataset. Then preprocessing: 70/15/15 split, ImageNet normalization. The two models are trained in parallel under identical conditions. Finally evaluation on the shared test set. The diagram also shows the key characteristics of each model side by side — you can already see the accuracy gap: 89.23% for SimpleCNN versus 96.10% for U-Net.""",

    14: """[~1:15]
SimpleCNN in numbers. The layer table shows that the dominant cost is the first fully connected layer: 8,192 × 512 = ~4.2 million parameters out of 4.5 million total. This is why Dropout at that layer is so important — without it, that many free parameters on 18,900 training images would overfit quickly. The convolutional layers themselves are lightweight: the three conv blocks together account for fewer than 100K parameters. The model fits in 17.2 MB and trains to its best checkpoint in just 11 epochs.""",

    15: """[~1:15]
U-Net Classifier in numbers. The parameter count is distributed much more evenly across the network. The encoder bottleneck — the 256→512 DoubleConv — accounts for ~3.5M parameters. The decoder adds another ~2M in upsampling and DoubleConv blocks. The classification head — AdaptiveAvgPool plus a single linear layer of 64×10 = 640 weights — is almost negligible. Total ~8.1M, model size 30.9 MB — 1.8× larger than SimpleCNN. Best epoch 24, test accuracy 96.10%.""",

    16: """[~1:30]
Both models share identical training conditions — this is essential for a fair comparison. Adam optimizer with learning rate 0.001 proved the most stable; lr=0.01 caused oscillation in early epochs. CrossEntropyLoss is the standard multi-class objective. Batch size 64 balances gradient noise and memory. The early stopping mechanism checks validation loss every epoch: if it improves by more than min_delta, save a checkpoint and reset the patience counter; if no improvement for 5 consecutive epochs, stop training and load the best saved checkpoint. The min_delta differs: 0.001 for SimpleCNN (noisier curve) and 0.0 for U-Net (smoother curve). Early stopping saved 34 unnecessary epochs for SimpleCNN and 21 for U-Net.""",

    17: """[~1:00]
Our test protocol uses the 4,050 held-out images seen at no point during training or hyperparameter selection. We report four metrics: accuracy, weighted F1, weighted precision, and weighted recall. Weighted averaging weights each class's score by its frequency in the test set — correct for a mildly imbalanced dataset. We also inspect the full 10×10 confusion matrix, which reveals which class pairs are most often confused. The checkpoint used is always the one with the lowest validation loss — not the last epoch — because validation loss starts degrading before accuracy does, giving an earlier and more reliable signal.""",

    18: """[~1:45]
SimpleCNN results. The training curves on the left show loss and accuracy over 16 epochs, with the best checkpoint at epoch 11 marked by a vertical line. Validation loss at best epoch: 0.3105. On the test set: accuracy 89.23%, weighted F1 0.8930. Looking at per-class results, the model handles visually distinctive classes well — Forest 0.954, SeaLake 0.948, Residential 0.939. The hardest classes are HerbaceousVegetation at 0.821, PermanentCrop at 0.833, and River at 0.836. These failures are predictable: at 10 meters per pixel, a river is a thin strip easily confused with a highway or surrounding fields; herbaceous vegetation and pastures look nearly identical.""",

    19: """[~1:45]
U-Net Classifier results. Best validation loss: 0.1116 at epoch 24. On the test set: accuracy 96.10%, weighted F1 0.9610 — a 6.87 percentage point improvement over SimpleCNN. What is most striking is that the gains are consistent across all classes. River jumps from 0.836 to 0.951 — a gain of +0.115 from one architectural change. PermanentCrop: +0.107. HerbaceousVegetation: +0.118. Even the already-strong classes improve: Highway from 0.848 to 0.978, Residential from 0.939 to 0.986. Why? Skip connections allow the network to simultaneously access fine spatial detail from shallow encoder layers and abstract semantic features from deep layers. For distinguishing visually similar land types, this dual access to local and global context is the key advantage.""",

    20: """[~1:30]
The comparison table summarizes everything. Accuracy improves 6.87 pp. Weighted F1 up 0.068. Best validation loss drops 64%. U-Net takes twice as many epochs and is 1.8× larger — but the worst-class F1 improves by 0.118. The bar charts confirm: every single metric improves, and every single class improves. The cost — 1.8× more parameters, 2× more training epochs, 13 extra MB of storage — is clearly justified. In any real deployment scenario where accuracy matters, U-Net is the better choice.""",

    21: """[~1:15]
To summarize: we designed, trained, and evaluated two architectures. SimpleCNN gives 89.23% accuracy with 4.5M parameters in 11 epochs — a strong and practical lightweight baseline. U-Net Classifier gives 96.10% with 8.1M parameters in 24 epochs — a systematic +6.87 pp across all classes. The central lesson: skip connections and hierarchical feature reuse, the defining properties of U-Net, are highly effective for satellite imagery because they allow the model to exploit both local texture and global context simultaneously. From an engineering standpoint, fixed seed, Docker, and the modular pipeline ensure full reproducibility. Looking ahead, the most promising directions are using all 13 Sentinel-2 bands, adding data augmentation, and fine-tuning a pre-trained backbone.""",

    22: """[~0:30]
Here are the references we cited throughout the presentation. The key papers are Helber et al. for EuroSAT, Ronneberger et al. for U-Net, and He et al. for ResNet. Thank you very much for your attention — we are happy to take any questions.""",
}

# ─────────────────────────────────────────────────────────────────────────────
# ASSEMBLE
# ─────────────────────────────────────────────────────────────────────────────

def add_notes(slide, text):
    """Attach speaker notes XML to a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text

def build():
    prs = make_prs()

    builders = [
        slide_01_title,
        slide_02_outline,
        slide_03_aim,
        slide_04_why,
        slide_05_dataset_overview,
        slide_06_classes,
        slide_07_samples,
        slide_08_prepare_folders,
        slide_09_prepare_split,
        slide_10_reading,
        slide_11_simplecnn,
        slide_12_unet,
        slide_13_arch_compare,
        slide_14_cnn_chars,
        slide_15_unet_chars,
        slide_16_training,
        slide_17_testing,
        slide_18_results_cnn,
        slide_19_results_unet,
        slide_20_comparison,
        slide_21_conclusion,
        slide_22_references,
    ]

    for i, builder in enumerate(builders, start=1):
        sl = builder(prs)
        if i in NOTES:
            add_notes(sl, NOTES[i])
        print(f"  Slide {i:2d} done — {builder.__name__}")

    out_pptx = "Land_Cover_Classification_NEW.pptx"
    prs.save(out_pptx)
    print(f"\nSaved: {out_pptx}")

    # Write speaker notes as standalone txt
    out_txt = "Speaker_Notes_NEW.txt"
    lines = [
        "=" * 80,
        "SPEAKER NOTES — Land Cover Classification from Satellite Imagery",
        "Estimated duration: 20–25 minutes",
        "=" * 80,
        "",
        "Tips:",
        "  • ~130 words per minute pace",
        "  • Each slide note is labeled [mm:ss]",
        "  • Total: ~22 minutes",
        "",
    ]
    slide_titles = [
        "SLIDE 1 — Title",
        "SLIDE 2 — Outline",
        "SLIDE 3 — Aim of the Work",
        "SLIDE 4 — Related Work & Limitations",
        "SLIDE 5 — EuroSAT Dataset Overview",
        "SLIDE 6 — 10 Land Cover Classes",
        "SLIDE 7 — Sample Patches",
        "SLIDE 8 — Folder Structure",
        "SLIDE 9 — Train/Val/Test Split",
        "SLIDE 10 — Reading the Dataset",
        "SLIDE 11 — SimpleCNN Architecture",
        "SLIDE 12 — U-Net Classifier Architecture",
        "SLIDE 13 — Architecture Overview (Pipeline)",
        "SLIDE 14 — SimpleCNN Characteristics",
        "SLIDE 15 — U-Net Characteristics",
        "SLIDE 16 — Training Techniques",
        "SLIDE 17 — Testing Protocol",
        "SLIDE 18 — Results: SimpleCNN",
        "SLIDE 19 — Results: U-Net Classifier",
        "SLIDE 20 — Model Comparison",
        "SLIDE 21 — Conclusion",
        "SLIDE 22 — References",
    ]
    for i, title in enumerate(slide_titles, start=1):
        lines.append("-" * 80)
        lines.append(title)
        lines.append("-" * 80)
        lines.append(NOTES.get(i, "(no notes)"))
        lines.append("")

    with open(out_txt, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))
    print(f"Saved: {out_txt}")

if __name__ == "__main__":
    build()
